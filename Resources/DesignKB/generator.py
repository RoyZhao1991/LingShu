#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
灵枢 自进化 PPT 生成器(设计系统驱动)。
用法: python3 generator.py <slides.json> <out.pptx> [designkb_dir]

读 DesignKB(palettes/typography),吃富 slides.json(每页带 layout),按版式原型组合排版——
封面/章节/要点(带图标)/大数字/左右图文/满版图/两栏对比/时间线/引言/图表/目录/收尾。
图片走 PIL 等比裁切填充(不拉伸);图标取 DesignKB/icons;图表用 python-pptx 原生 chart。
设计是数据驱动的:换 palette / 调 layout 都不动这段代码 —— 这是自进化的基础。
"""
import sys, json, os, tempfile, zipfile, re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

EMU_PER_IN = 914400

def _hex(c):
    c = c.lstrip('#')
    return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))

def load_json(path, default):
    try:
        with open(path, encoding='utf-8') as f:
            return json.load(f)
    except Exception:
        return default

# ---------- 模板主题提取(让"以模板为底"真正改变观感) ----------
# 给了 template 时,不再只继承母版——从模板 theme1.xml 提取配色 + 字体,推导成生成器用的
# palette,让全篇采用模板自身的品牌色/字体(这才是"参考模板"的意义)。提取失败回退 DesignKB 配色。
_CJK_FONTS = ('pingfang', 'yahei', 'source han', 'noto sans cjk', 'noto serif cjk',
              'simhei', 'simsun', 'songti', 'heiti', 'jhenghei', 'hiragino', 'kai', 'fangsong')

def _lum(hexc):
    c = hexc.lstrip('#')
    r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
    return (0.299 * r + 0.587 * g + 0.114 * b) / 255.0

def _mix(a, b, t):
    """按 t∈[0,1] 在两个 #RRGGBB 间线性插值。"""
    a, b = a.lstrip('#'), b.lstrip('#')
    return '#' + ''.join('%02X' % round(int(a[i:i+2], 16) * (1 - t) + int(b[i:i+2], 16) * t) for i in (0, 2, 4))

def extract_template_theme(path):
    """从模板 .pptx 的 theme1.xml 抽 clrScheme + fontScheme,推导 {palette, fonts}。失败/无主题返回 None。"""
    try:
        with zipfile.ZipFile(path) as z:
            name = next((n for n in z.namelist() if re.match(r'ppt/theme/theme\d+\.xml$', n)), None)
            if not name:
                return None
            xml = z.read(name).decode('utf-8', 'ignore')
    except Exception:
        return None

    def clr(tag):
        m = re.search(r'<a:%s>\s*<a:(?:srgbClr val="([0-9A-Fa-f]{6})"|sysClr val="\w+" lastClr="([0-9A-Fa-f]{6})")' % tag, xml)
        return ('#' + (m.group(1) or m.group(2)).upper()) if m else None

    dk1, lt1 = clr('dk1') or '#000000', clr('lt1') or '#FFFFFF'
    dk2, lt2 = clr('dk2'), clr('lt2')
    a1, a2 = clr('accent1'), clr('accent2')
    if not a1:
        return None  # 没有主强调色 → 主题信息不足,回退 DesignKB
    # 标准 office 约定:lt1↔背景、dk1↔正文。按背景亮度判明暗主题。
    bg, ink = lt1, dk1
    if _lum(bg) < 0.5:
        mode = 'dark'; surface = _mix(bg, '#FFFFFF', 0.08); muted = _mix(ink, bg, 0.45)
    else:
        mode = 'light'; surface = _mix(bg, a1, 0.05); muted = _mix(ink, bg, 0.5)
    pal = {'bg': bg, 'surface': surface, 'ink': ink, 'muted': muted,
           'accent': a1, 'accent2': a2 or dk2 or a1, 'mode': mode}

    # 字体:仅当能渲染中文(ea 或本身是 CJK 字体)才采用,否则保留 CJK 安全默认(避免中文掉字)。
    def pick_font(block_tag):
        blk = re.search(r'<a:%sFont>(.*?)</a:%sFont>' % (block_tag, block_tag), xml, re.S)
        if not blk:
            return None
        body = blk.group(1)
        ea = (re.search(r'<a:ea typeface="([^"]*)"', body) or [None, ''])[1]
        latin = (re.search(r'<a:latin typeface="([^"]*)"', body) or [None, ''])[1]
        for cand in (ea, latin):
            if cand and any(k in cand.lower() for k in _CJK_FONTS):
                return cand
        return None

    fonts = {}
    tf, bf = pick_font('major'), pick_font('minor')
    if tf:
        fonts['title'] = tf
    if bf:
        fonts['body'] = bf
    return {'palette': pal, 'fonts': fonts}

# ---------- DesignKB ----------
SRC = sys.argv[1] if len(sys.argv) > 1 else 'slides.json'
OUT = sys.argv[2] if len(sys.argv) > 2 else '演示.pptx'
KB = sys.argv[3] if len(sys.argv) > 3 else os.path.dirname(os.path.abspath(__file__))

palettes = {p['id']: p for p in load_json(os.path.join(KB, 'palettes.json'), {}).get('palettes', [])}
typo = load_json(os.path.join(KB, 'typography.json'), {})
SCALE = typo.get('scale', {})
GRID = typo.get('grid', {})
ICON_DIR = os.path.join(KB, 'icons', 'lucide')

data = load_json(SRC, {})
slides = data.get('slides', data if isinstance(data, list) else [])

theme = data.get('theme', 'midnight')
pal_id = theme.get('palette') if isinstance(theme, dict) else theme
PAL = palettes.get(pal_id) or (list(palettes.values())[0] if palettes else {
    "bg": "#0B1220", "surface": "#131C2E", "ink": "#F2FFFD", "muted": "#9FB8B4", "accent": "#25F4E4", "accent2": "#7C5CFF", "mode": "dark"})
TITLE_FONT = (typo.get('pairings', [{}])[0]).get('title_font', 'PingFang SC')
BODY_FONT = (typo.get('pairings', [{}])[0]).get('body_font', 'PingFang SC')

# 模板主题接管=**显式 opt-in**(默认尊重 DesignKB 配色——它本身就是打磨好的专业主题,通常比联网套来的
# 通用 .pptx 好看;实测自动套用通用 Office 模板会把深色专业风变成寡淡白底,反而更差)。
# 只有显式要求(theme=="template" 或 {"palette_source":"template"})才用模板自身配色/字体,通常用于
# "以用户自带品牌模板为准"。普通 theme(如 midnight)一律按 DesignKB 配色,不被模板覆盖。
_tmpl_path = data.get('template')
_want_tmpl_theme = (isinstance(theme, str) and theme.strip().lower() == 'template') \
    or (isinstance(theme, dict) and theme.get('palette_source') == 'template')
_tmpl_theme = extract_template_theme(_tmpl_path) if (_want_tmpl_theme and _tmpl_path and os.path.exists(_tmpl_path)) else None
THEME_SOURCE = pal_id or 'default'
if _tmpl_theme:
    PAL = _tmpl_theme['palette']
    TITLE_FONT = _tmpl_theme['fonts'].get('title', TITLE_FONT)
    BODY_FONT = _tmpl_theme['fonts'].get('body', BODY_FONT)
    THEME_SOURCE = 'template'

BG, SURFACE, INK, MUTED = _hex(PAL['bg']), _hex(PAL['surface']), _hex(PAL['ink']), _hex(PAL['muted'])
ACCENT, ACCENT2 = _hex(PAL['accent']), _hex(PAL.get('accent2', PAL['accent']))
DARK_MODE = PAL.get('mode', 'dark') == 'dark'

SW = GRID.get('slide_w_in', 13.333); SH = GRID.get('slide_h_in', 7.5)
M = GRID.get('margin_in', 0.9)

# 模板底模式:给了 template 且文件存在 → 以它为底(继承母版,主题配色/字体已在上面提取并接管),
# 清掉它自带的样例页;否则从空白起手。这样 acquire_resource 找来的专业模板能真正被"参考/做底"。
if _tmpl_path and os.path.exists(_tmpl_path):
    try:
        prs = Presentation(_tmpl_path)
        _ids = prs.slides._sldIdLst
        _RID = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        for _s in list(_ids):
            # 既要从 sldIdLst 摘掉,也要 drop 关系——否则样例页 part 残留,保存时与新页 slide1.xml 撞名(损坏文件)。
            _rid = _s.get(_RID)
            if _rid:
                try: prs.part.drop_rel(_rid)
                except Exception: pass
            _ids.remove(_s)   # 清掉模板自带样例页,只留母版/主题
    except Exception:
        prs = Presentation()
else:
    prs = Presentation()
prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

def sz(key, fallback): return Pt(SCALE.get(key, fallback))

def add_slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill; f.solid(); f.fore_color.rgb = bg
    return s

def rect(s, l, t, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=8):
    """runs: list of (string, size_pt, color, bold, font). 多段=多 paragraph。"""
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, item in enumerate(runs):
        txt, size, color, bold, font = (list(item) + [None])[:5]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        r = p.add_run(); r.text = str(txt)
        r.font.size = size; r.font.bold = bold; r.font.color.rgb = color
        r.font.name = font or BODY_FONT
    return box

def title_pt(t):
    n = len(t or '')
    return sz('title', 32) if n <= 16 else (Pt(28) if n <= 26 else Pt(24))

def _est_lines(text, width_in, font_pt):
    """粗估文本在给定宽度(英寸)/字号(pt)下换行行数(CJK≈字号宽,ASCII≈0.55×)。用于动态排版防重叠。"""
    if not text:
        return 1
    width_pt = max(1.0, width_in * 72.0)
    w = 0.0
    for ch in str(text):
        w += font_pt * (0.55 if ord(ch) < 0x2E80 else 1.0)
    return max(1, int((w + width_pt - 0.001) // width_pt))

def place_image(s, path, l, t, w, h):
    """等比裁切填充目标区域(cover),不拉伸变形。失败→返回 False。"""
    if not path or not os.path.exists(path):
        return False
    try:
        from PIL import Image
        im = Image.open(path).convert('RGB')
        iw, ih = im.size; tr = w / h; ir = iw / ih
        if ir > tr:
            nw = int(ih * tr); x = (iw - nw) // 2; im = im.crop((x, 0, x + nw, ih))
        else:
            nh = int(iw / tr); y = (ih - nh) // 2; im = im.crop((0, y, iw, y + nh))
        im = im.resize((int(w * 200), int(h * 200)))
        tmp = tempfile.mktemp(suffix='.jpg'); im.save(tmp, 'JPEG', quality=88)
        s.shapes.add_picture(tmp, Inches(l), Inches(t), Inches(w), Inches(h))
        return True
    except Exception:
        try:
            s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
            return True
        except Exception:
            return False

def place_icon(s, name, l, t, size=0.42):
    """放图标:**只用与背景对比的变体**(深色主题→白,浅色主题→深)。对比变体缺失就返回 False
    交给调用方画亮色圆点兜底——绝不退用同色调图标(那会在深色底变成看不见的黑图标,曾被用户指出)。"""
    if not name:
        return False
    variant = 'white' if DARK_MODE else 'dark'
    p = os.path.join(ICON_DIR, f"{name}-{variant}.png")
    if os.path.exists(p):
        try:
            s.shapes.add_picture(p, Inches(l), Inches(t), Inches(size), Inches(size)); return True
        except Exception:
            return False
    return False   # 无对比变体 → 不硬塞同色图标,让调用方用 ACCENT 圆点(始终可见)

def kicker(s, txt, l=M, t=0.62):
    rect(s, l, t + 0.07, 0.32, 0.06, fill=ACCENT)
    text(s, l + 0.42, t, 6, 0.4, [(txt, Pt(13), ACCENT, True, BODY_FONT)])

def footer(s, idx, total, deck):
    text(s, M, SH - 0.5, 6, 0.3, [(deck or '', Pt(10), MUTED, False, BODY_FONT)])
    text(s, SW - M - 2, SH - 0.5, 2, 0.3, [(f"{idx:02d} / {total:02d}", Pt(10), MUTED, False, BODY_FONT)], align=PP_ALIGN.RIGHT)

# ---------- 版式原型 ----------
def L_cover(s, d, **k):
    rect(s, 0, 0, SW, SH, fill=BG)
    img = place_image(s, d.get('image'), SW * 0.52, 0, SW * 0.48, SH)
    if img:
        rect(s, SW * 0.52, 0, SW * 0.04, SH, fill=BG)  # 接缝过渡
        text_w = SW * 0.52 - M - 0.35   # 文字严格限制在左栏,绝不压到右侧图
    else:
        rect(s, SW - 3.2, 0, 3.2, SH, fill=SURFACE)
        rect(s, SW - 3.2, 0, 0.12, SH, fill=ACCENT)
        text_w = SW - 3.2 - M - 0.4
    rect(s, M, SH * 0.3, 0.55, 0.1, fill=ACCENT)
    # 标题字号按长度自适应,确保在左栏宽度内一两行放下、不溢出不触边。
    title = str(d.get('title', '标题'))
    base = SCALE.get('cover_title', 54)
    tsize = base if len(title) <= 8 else (44 if len(title) <= 13 else (36 if len(title) <= 19 else 30))
    # 标题占高随字号/行数估算,副标题紧随其后(动态推进 y,杜绝长副标题换行压住 tagline)。
    title_lines = max(1, _est_lines(title, text_w, tsize))
    text(s, M, SH * 0.34, text_w, 0.46 * tsize / 32 * title_lines + 0.4, [(title, Pt(tsize), INK, True, TITLE_FONT)], anchor=MSO_ANCHOR.TOP)
    y = SH * 0.34 + 0.5 * (tsize / 32.0) * title_lines + 0.5
    if d.get('subtitle'):
        sub = str(d['subtitle'])
        ssize = 22 if len(sub) <= 14 else (18 if len(sub) <= 26 else 16)
        slines = max(1, _est_lines(sub, text_w, ssize))
        text(s, M, y, text_w, 0.34 * slines + 0.2, [(sub, Pt(ssize), ACCENT, True, BODY_FONT)])
        y += 0.32 * slines + 0.28
    if d.get('tagline'):
        text(s, M, y, text_w, 1.2, [(d['tagline'], Pt(15), MUTED, False, BODY_FONT)])

def L_section(s, d, **k):
    rect(s, 0, 0, SW, SH, fill=SURFACE)
    rect(s, 0, 0, 0.18, SH, fill=ACCENT)
    text(s, M, SH * 0.28, 4, 1.6, [(str(d.get('index', '')), Pt(96), ACCENT, True, TITLE_FONT)])
    text(s, M, SH * 0.55, SW - 2 * M, 1.6, [(d.get('title', ''), sz('section_title', 40), INK, True, TITLE_FONT)])
    if d.get('subtitle'):
        text(s, M, SH * 0.55 + 1.0, SW - 2 * M, 0.8, [(d['subtitle'], Pt(16), MUTED, False, BODY_FONT)])

def _bullet_block(s, bullets, icons, l, t, w, h):
    n = max(1, len(bullets)); row = min(0.92, h / n)
    for i, b in enumerate(bullets):
        y = t + i * row
        ic = icons[i] if icons and i < len(icons) else None
        placed = place_icon(s, ic, l, y + 0.04, size=0.4)
        if not placed:
            rect(s, l + 0.05, y + 0.16, 0.16, 0.16, fill=ACCENT, shape=MSO_SHAPE.OVAL)
        text(s, l + 0.62, y, w - 0.62, row, [(str(b), sz('bullet', 18), INK if DARK_MODE else INK, False, BODY_FONT)], anchor=MSO_ANCHOR.MIDDLE, space=2)

def L_bullets(s, d, **k):
    kicker(s, d.get('kicker', '要点'))
    text(s, M, 1.15, SW - 2 * M, 1.2, [(d.get('title', ''), title_pt(d.get('title')), INK, True, TITLE_FONT)])
    _bullet_block(s, d.get('bullets', []), d.get('icons'), M, 2.5, SW - 2 * M, SH - 3.2)

def L_bignumber(s, d, **k):
    kicker(s, d.get('kicker', '关键指标'))
    text(s, M, 1.5, SW - 2 * M, 2.2, [(str(d.get('number', '')), sz('bignum', 120), ACCENT, True, TITLE_FONT)])
    text(s, M, 3.9, SW - 2 * M, 0.8, [(d.get('label', ''), sz('bignum_label', 22), INK, True, BODY_FONT)])
    if d.get('title') or d.get('desc'):
        text(s, M, 4.7, SW - 2 * M, 1.6, [(d.get('title') or d.get('desc'), Pt(16), MUTED, False, BODY_FONT)])

def _image_text(s, d, image_left):
    iw = SW * 0.46
    ix = 0 if image_left else SW - iw
    tx = iw + 0.7 if image_left else M
    ok = place_image(s, d.get('image'), ix, 0, iw, SH)
    if not ok:
        rect(s, ix, 0, iw, SH, fill=SURFACE); rect(s, ix, 0, iw, 0.12, fill=ACCENT)
    kicker(s, d.get('kicker', '要点'), l=tx)
    text(s, tx, 1.15, SW - iw - M - 0.7, 1.4, [(d.get('title', ''), title_pt(d.get('title')), INK, True, TITLE_FONT)])
    _bullet_block(s, d.get('bullets', []), d.get('icons'), tx, 2.6, SW - iw - M - 0.7, SH - 3.3)

def L_image_right(s, d, **k): _image_text(s, d, image_left=False)
def L_image_left(s, d, **k): _image_text(s, d, image_left=True)

def L_image_full(s, d, **k):
    ok = place_image(s, d.get('image'), 0, 0, SW, SH)
    if not ok:
        rect(s, 0, 0, SW, SH, fill=SURFACE)
    rect(s, 0, SH - 2.6, SW, 2.6, fill=BG)  # 底部遮罩
    rect(s, M, SH - 2.1, 0.55, 0.1, fill=ACCENT)
    text(s, M, SH - 1.95, SW - 2 * M, 1.4, [(d.get('title', ''), Pt(34), _hex('#FFFFFF'), True, TITLE_FONT)])

def _card(s, l, t, w, h, head, bullets):
    rect(s, l, t, w, h, fill=SURFACE)
    rect(s, l, t, w, 0.1, fill=ACCENT)
    text(s, l + 0.35, t + 0.35, w - 0.7, 0.7, [(head, Pt(20), INK, True, TITLE_FONT)])
    text(s, l + 0.35, t + 1.15, w - 0.7, h - 1.4,
         [(("• " + str(b)), sz('bullet', 18), MUTED, False, BODY_FONT) for b in (bullets or [])], space=8)

def L_twocol(s, d, **k):
    kicker(s, d.get('kicker', '对比'))
    text(s, M, 1.15, SW - 2 * M, 1.2, [(d.get('title', ''), title_pt(d.get('title')), INK, True, TITLE_FONT)])
    cw = (SW - 2 * M - 0.6) / 2
    left = d.get('left', {}); right = d.get('right', {})
    _card(s, M, 2.7, cw, SH - 3.4, left.get('heading', '左'), left.get('bullets', []))
    _card(s, M + cw + 0.6, 2.7, cw, SH - 3.4, right.get('heading', '右'), right.get('bullets', []))

def L_timeline(s, d, **k):
    kicker(s, d.get('kicker', '路线'))
    text(s, M, 1.15, SW - 2 * M, 1.2, [(d.get('title', ''), title_pt(d.get('title')), INK, True, TITLE_FONT)])
    steps = d.get('steps', [])[:5]; n = max(1, len(steps))
    cw = (SW - 2 * M) / n; y = 3.6
    rect(s, M, y + 0.18, SW - 2 * M, 0.04, fill=MUTED)
    for i, st in enumerate(steps):
        cx = M + i * cw
        rect(s, cx + cw / 2 - 0.16, y + 0.02, 0.32, 0.32, fill=ACCENT, shape=MSO_SHAPE.OVAL)
        text(s, cx, y - 0.7, cw, 0.6, [(st.get('label', f'阶段{i+1}'), Pt(16), INK, True, TITLE_FONT)], align=PP_ALIGN.CENTER)
        text(s, cx, y + 0.55, cw, 1.6, [(st.get('desc', ''), Pt(13), MUTED, False, BODY_FONT)], align=PP_ALIGN.CENTER)

def L_quote(s, d, **k):
    rect(s, 0, 0, SW, SH, fill=SURFACE)
    text(s, M + 0.2, SH * 0.18, 1.2, 1.2, [("“", Pt(96), ACCENT, True, TITLE_FONT)])
    text(s, M, SH * 0.33, SW - 2 * M, 2.6, [(d.get('quote', ''), sz('quote', 30), INK, True, TITLE_FONT)], anchor=MSO_ANCHOR.TOP)
    if d.get('attrib'):
        text(s, M, SH * 0.78, SW - 2 * M, 0.6, [("— " + d['attrib'], Pt(16), ACCENT, True, BODY_FONT)])

def _strip_lead_num(t):
    """去掉条目开头已有的序号前缀(01 / 1. / 1、 / 1) 等),避免与版式自带编号重复(曾"01 01 定位")。"""
    s = str(t)
    return re.sub(r'^\s*\d{1,2}(?:[.,、)\]：:．-]\s*|\s+)', '', s).strip() or s

def L_agenda(s, d, **k):
    kicker(s, '目录')
    text(s, M, 1.15, SW - 2 * M, 1.2, [(d.get('title', '目录'), title_pt(d.get('title') or '目录'), INK, True, TITLE_FONT)])
    items = d.get('items', []); y = 2.6
    for i, it in enumerate(items[:6]):
        text(s, M, y, 1.0, 0.6, [(f"{i+1:02d}", Pt(22), ACCENT, True, TITLE_FONT)])
        text(s, M + 1.0, y + 0.03, SW - 2 * M - 1.0, 0.6, [(_strip_lead_num(it), Pt(19), INK, False, BODY_FONT)], anchor=MSO_ANCHOR.MIDDLE)
        y += 0.72

def L_chart(s, d, **k):
    kicker(s, d.get('kicker', '数据'))
    text(s, M, 1.15, SW - 2 * M, 1.0, [(d.get('title', ''), title_pt(d.get('title')), INK, True, TITLE_FONT)])
    ch = d.get('chart', {})
    try:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        cd = CategoryChartData(); cd.categories = ch.get('categories', [])
        for ser in ch.get('series', []):
            cd.add_series(ser.get('name', ''), ser.get('values', []))
        kind = {'bar': XL_CHART_TYPE.COLUMN_CLUSTERED, 'line': XL_CHART_TYPE.LINE_MARKERS,
                'pie': XL_CHART_TYPE.PIE}.get(ch.get('type', 'bar'), XL_CHART_TYPE.COLUMN_CLUSTERED)
        gf = s.shapes.add_chart(kind, Inches(M), Inches(2.5), Inches(SW - 2 * M), Inches(SH - 3.2), cd)
        chart = gf.chart; chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        # 图表文字随主题(深色主题用浅墨色),保证轴/图例可读。
        try:
            chart.font.size = Pt(12); chart.font.color.rgb = INK
            chart.legend.font.color.rgb = MUTED
        except Exception:
            pass
    except Exception as e:
        text(s, M, 3.0, SW - 2 * M, 2, [(f"(图表数据缺失或不可用:{e})", Pt(14), MUTED, False, BODY_FONT)])

def L_closing(s, d, **k):
    rect(s, 0, 0, SW, SH, fill=BG)
    rect(s, 0, SH - 1.4, SW, 1.4, fill=ACCENT)
    text(s, M, SH * 0.32, SW - 2 * M, 1.8, [(d.get('title', '谢谢'), Pt(46), INK, True, TITLE_FONT)])
    if d.get('bullets'):
        text(s, M, SH * 0.32 + 1.5, SW - 2 * M, 1.6,
             [(("→ " + str(b)), Pt(18), MUTED, False, BODY_FONT) for b in d['bullets']], space=6)
    contact = _clean_contact(d.get('contact'))
    if contact:
        text(s, M, SH - 1.15, SW - 2 * M, 0.8, [(contact, Pt(16), BG, True, BODY_FONT)])

def _clean_contact(c):
    """收尾页联系信息防泄露:绝不渲染绝对文件路径/工作目录(曾把 /Users/example/app 当 contact 印上去)。"""
    c = str(c or '').strip()
    if not c:
        return ''
    if '/Users/' in c or c.startswith('/') or '\\' in c or '\\\\' in c:
        kept = [p for p in re.split(r'\s+', c) if not (p.startswith('/') or '/Users/' in p or '\\' in p)]
        return ' '.join(kept).strip()
    return c

LAYOUTS = {
    'cover': L_cover, 'section': L_section, 'bullets': L_bullets, 'bignumber': L_bignumber,
    'image-right': L_image_right, 'image-left': L_image_left, 'image-full': L_image_full,
    'twocol': L_twocol, 'timeline': L_timeline, 'quote': L_quote, 'agenda': L_agenda,
    'chart': L_chart, 'closing': L_closing,
}

deck_title = data.get('title', '')
total = len(slides)
used = []
for i, d in enumerate(slides, 1):
    lay = d.get('layout', 'bullets')
    fn = LAYOUTS.get(lay, L_bullets)
    used.append(lay if lay in LAYOUTS else 'bullets')
    s = add_slide()
    try:
        fn(s, d)
    except Exception as e:
        text(s, M, 1.2, SW - 2 * M, 2, [(d.get('title', ''), Pt(28), INK, True, TITLE_FONT), (f"(版式渲染降级:{e})", Pt(12), MUTED, False, BODY_FONT)])
    if lay not in ('cover', 'section', 'quote', 'image-full', 'closing'):
        footer(s, i, total, deck_title)

prs.save(OUT)
print(f"OK pages={total} theme={THEME_SOURCE} palette={pal_id} layouts={','.join(used)}")
